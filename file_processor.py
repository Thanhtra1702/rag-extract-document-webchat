import csv
import io
import json
import os

from langchain_core.documents import Document
from langchain_classic.text_splitter import RecursiveCharacterTextSplitter


# ---------------------------------------------------------------------------
# Text extraction helpers (one per format)
# ---------------------------------------------------------------------------

def _extract_pdf(file_bytes: bytes) -> str:
    try:
        from PyPDF2 import PdfReader
    except ImportError as exc:
        raise RuntimeError("PDF support requires PyPDF2. Install it in the active Python environment.") from exc
    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text.strip())
    return "\n\n".join(pages)


def _extract_docx(file_bytes: bytes) -> str:
    from docx import Document as DocxDocument
    doc = DocxDocument(io.BytesIO(file_bytes))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_txt(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return file_bytes.decode("latin-1")


def _extract_csv(file_bytes: bytes) -> str:
    text = _extract_txt(file_bytes)
    reader = csv.reader(io.StringIO(text))
    rows = []
    for row in reader:
        rows.append(" | ".join(row))
    return "\n".join(rows)


def _strip_html(text: str) -> str:
    """Remove HTML tags from a string."""
    import re
    return re.sub(r'<[^>]+>', '', text).strip()


def _flatten_json_value(data, prefix: str = "") -> list[str]:
    """Recursively flatten JSON data into readable key-value lines."""
    lines = []
    if isinstance(data, dict):
        for key, value in data.items():
            new_prefix = f"{prefix}{key}" if not prefix else f"{prefix} > {key}"
            if isinstance(value, (dict, list)):
                lines.extend(_flatten_json_value(value, new_prefix))
            else:
                clean_val = _strip_html(str(value)) if isinstance(value, str) else str(value)
                if clean_val:
                    lines.append(f"{new_prefix}: {clean_val}")
    elif isinstance(data, list):
        for i, item in enumerate(data):
            if isinstance(item, dict):
                # For product-like objects, flatten each as a block
                item_lines = _flatten_json_value(item, "")
                if item_lines:
                    lines.append(f"\n--- Item {i+1} ---")
                    lines.extend(item_lines)
            else:
                clean_val = _strip_html(str(item)) if isinstance(item, str) else str(item)
                if clean_val:
                    lines.append(f"{prefix}[{i}]: {clean_val}")
    else:
        clean_val = _strip_html(str(data)) if isinstance(data, str) else str(data)
        if clean_val:
            lines.append(f"{prefix}: {clean_val}")
    return lines


def _extract_json(file_bytes: bytes) -> str:
    # Thử giải mã bằng nhiều encoding thông dụng
    text = None
    for enc in ["utf-8", "utf-8-sig", "utf-16", "latin-1", "cp1258"]:
        try:
            text = file_bytes.decode(enc)
            break
        except UnicodeDecodeError:
            continue
            
    if text is None:
        raise ValueError("Could not decode the JSON file. Please save it as UTF-8.")
        
    try:
        data = json.loads(text)
        # Flatten JSON into readable plain text (strip HTML, key-value pairs)
        lines = _flatten_json_value(data)
        return "\n".join(lines)
    except json.JSONDecodeError:
        # Nếu file JSON bị lỗi cú pháp nhẹ, fallback đọc trực tiếp như file text để tránh lỗi tải lên
        return text



def _extract_xlsx(file_bytes: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"--- Sheet: {sheet_name} ---")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)


def _extract_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            slides_text.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides_text)


# ---------------------------------------------------------------------------
# Format dispatcher
# ---------------------------------------------------------------------------

def _extract_image(file_bytes: bytes) -> str:
    """Images have no extractable text without OCR/VLM. Return a placeholder."""
    return "[Image file uploaded — visual content available in Document Parser view]"


EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".txt": _extract_txt,
    ".csv": _extract_csv,
    ".json": _extract_json,
    ".xlsx": _extract_xlsx,
    ".md": _extract_txt,
    ".pptx": _extract_pptx,
    ".png": _extract_image,
    ".jpg": _extract_image,
    ".jpeg": _extract_image,
    ".gif": _extract_image,
    ".webp": _extract_image,
    ".bmp": _extract_image,
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
SUPPORTED_EXTENSIONS = list(EXTRACTORS.keys())

# ---------------------------------------------------------------------------
# Chunking — RecursiveCharacterTextSplitter
# ---------------------------------------------------------------------------

_text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=500,
    chunk_overlap=50,
    length_function=len,
    separators=["\n\n", "\n", ". ", ", ", " ", ""],
)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def process_file(filename: str, file_bytes: bytes) -> list[Document]:
    """
    Process raw file bytes into chunked LangChain Documents.

    Parameters
    ----------
    filename : str
        Original filename (used to detect format and stored in metadata).
    file_bytes : bytes
        Raw file content.

    Returns
    -------
    list[Document]
        Chunked documents with metadata (source filename, chunk index).
    """
    ext = os.path.splitext(filename)[1].lower()

    if ext not in EXTRACTORS:
        raise ValueError(
            f"Unsupported file format: '{ext}'. "
            f"Supported: {', '.join(SUPPORTED_EXTENSIONS)}"
        )

    extractor = EXTRACTORS[ext]
    raw_text = extractor(file_bytes)

    if not raw_text or not raw_text.strip():
        return []

    # Use RecursiveCharacterTextSplitter for intelligent chunking
    chunks = _text_splitter.split_text(raw_text.strip())

    docs = []
    for i, chunk in enumerate(chunks):
        docs.append(
            Document(
                page_content=chunk,
                metadata={
                    "source": filename,
                    "chunk_index": i,
                    "total_chunks": len(chunks),
                },
            )
        )
    return docs
